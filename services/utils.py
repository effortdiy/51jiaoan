
import requests
import config
import os
from PIL import Image
import logging
import aiofiles


async def download_img(folder,url,file_type):
    """ 图片下载 """
    file_name = os.path.basename(url)
    response = http_get(url)
    if not os.path.exists(folder):
        os.makedirs(folder)

    dest_file_name = os.path.join(folder,file_name)
    if not dest_file_name.endswith(file_type):
        dest_file_name +=dest_file_name+"."+file_type
        
    async with aiofiles.open(dest_file_name, 'wb') as f:
        await f.write(response.content)

    return True

def http_get(url:str,query_param:str=None):
    if not url.startswith("http"):
        url = f"http:{url}"
    res = requests.get(url=url,params=query_param,headers=config.HEADER)
    assert res.status_code == 200,"请求失败"
    return res

def http_post(url:str,json:dict):
    if not url.startswith("http"):
        url = f"http:{url}"
    res = requests.post(url=url,json=json,headers=config.HEADER)
    assert res.status_code == 200,"请求失败"
    return res


def image2pdf(dest_file_name, src_image_file: list):
    from reportlab.pdfgen import canvas

    # 创建一个新的PDF文件
    pdf = canvas.Canvas(dest_file_name)

    # 将每个图像保存到PDF文件中
    for i, img in enumerate(src_image_file):
        image = Image.open(img)
        # width, height = image.size
        width, height = image.size

        # 创建一个新的页面并绘制图像
        if i > 0:
            pdf.showPage()

        pdf.setPageSize(image.size)
        pdf.drawImage(img, 0, 0, width=width, height=height)

    # 保存PDF文件
    pdf.save()


def image2ppt(dest_file_path, src_images: list):
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches

    # 创建PPT对象
    prs = Presentation()

    # 遍历PDF图像并创建幻灯片
    for i, image in enumerate(src_images):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # 使用标题幻灯片布局
        left = top = 0
        with Image.open(image) as _:
            width,heigt = _.size
            # 计算比例

            slide.shapes.add_picture(image, left, top,width=prs.slide_width,height=prs.slide_height)
    # 保存PPT
    prs.save(dest_file_path)

def trans_doc(src_doc_dir,dest_doc_dir):
    """
    图片转文件
    """
    logging.info("文档转换ing...")
    # 第一层是目录
    for dir_1 in os.listdir(src_doc_dir):
        dir_1 = os.path.join(src_doc_dir,dir_1)
        if not os.path.isdir(dir_1):
            continue

        
        # # 第二层是目录
        # for dir_2 in os.listdir(dir_1):
        #     dir_2 = os.path.join(dir_1,dir_2)
        #     if not os.path.isdir(dir_2):
        #         continue

        #     # 第三层是文件
        if len(os.listdir(dir_1)) == 0:
            continue

        images = []
        for img in os.listdir(dir_1):
            img = os.path.join(dir_1,img)
            assert os.path.isfile(img) ,f"存在:\n   {img},不是文件。无法转换!!!"
            assert os.path.splitext(img)[1] in ['.jpg','.png'] ,f"存在:\n {img},不是jpg、png文件。无法转换!!!"

            images.append(img)

        out_file_name = os.path.basename(os.path.splitext(dir_1)[0])
        out_file_path = "\\".join(os.path.dirname(dir_1).split("\\")[2:])
        out_file_name = os.path.join(dest_doc_dir,out_file_path,"z-"+out_file_name)
        
        if not os.path.exists(os.path.split(out_file_name)[0]):
            os.mkdir(os.path.split(out_file_name)[0])

        if dir_1.endswith("pptx") or dir_1.endswith("ppt"):
            image2ppt(out_file_name+".pptx",images)
        else:
            image2pdf(out_file_name+".pdf",images)
        logging.info(f"文档 {out_file_name} 转换成功!")

    logging.info("文档转换完成")


